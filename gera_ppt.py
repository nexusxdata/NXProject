from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL_ESCURO  = RGBColor(0x1A, 0x3A, 0x6B)
AZUL_MEDIO   = RGBColor(0x2B, 0x57, 0x9A)
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF7)
LARANJA      = RGBColor(0xE8, 0x7B, 0x00)
CINZA_ESCURO = RGBColor(0x2D, 0x2D, 0x2D)
CINZA_MEDIO  = RGBColor(0x55, 0x5F, 0x6E)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_BG     = RGBColor(0xF4, 0xF6, 0xFB)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape

def add_textbox(slide, l, t, w, h, text, size, bold=False, color=CINZA_ESCURO,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size, bold=False, color=CINZA_ESCURO,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header(slide, titulo):
    add_rect(slide, 0, 0, W, H, CINZA_BG)
    add_rect(slide, 0, 0, W, Inches(1.3), AZUL_ESCURO)
    add_rect(slide, 0, Inches(1.3), Inches(0.07), H - Inches(1.3), LARANJA)
    add_textbox(slide, Inches(0.5), Inches(0.22), Inches(12), Inches(0.8),
                titulo, 28, bold=True, color=BRANCO)

# ── Slide 1 — Capa ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, AZUL_ESCURO)
add_rect(slide, 0, Inches(4.6), W, Inches(0.08), LARANJA)
add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(1.4),
            "NXProject", 72, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
            "Planejamento inteligente integrado ao Azure DevOps",
            24, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
            "Da estimativa ao cronograma — sem abandonar o backlog técnico.",
            16, italic=True, color=RGBColor(0xB0, 0xC8, 0xF0), align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
            "Nexus XData Tecnologia  •  nexusxdata.com.br",
            11, color=RGBColor(0x7A, 0x9A, 0xC8), align=PP_ALIGN.CENTER)

# ── Slide 2 — O problema ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "O problema que o NXProject resolve")

dores = [
    ("📋  Backlog rico, cronograma inexistente",
     "O time técnico trabalha no Azure DevOps mas a gestão não consegue\nvisualizar datas, dependências e alocação de forma consolidada."),
    ("⚠️  Retrabalho e planilhas paralelas",
     "Gestores mantêm planilhas Excel manualmente, fora de sincronia com o\nque o time realmente está fazendo no backlog."),
    ("🔗  Falta de rastreabilidade entre planejamento e execução",
     "Não há conexão direta entre a estimativa do gestor e o work item\ndo desenvolvedor — qualquer mudança exige atualização dupla."),
    ("📅  Dificuldade de replanejar rapidamente",
     "Quando uma sprint atrasa, recalcular todas as dependências e datas\nmanualmente é demorado e sujeito a erros."),
]
for i, (titulo, corpo) in enumerate(dores):
    col = i % 2; row = i // 2
    l = Inches(0.5 + col * 6.4); t = Inches(1.7 + row * 2.5)
    w = Inches(6.0); h = Inches(2.2)
    add_rect(slide, l, t, w, h, BRANCO)
    tb = slide.shapes.add_textbox(l+Inches(0.18), t+Inches(0.15), w-Inches(0.3), h-Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = titulo
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = AZUL_ESCURO
    add_para(tf, corpo, 11, color=CINZA_MEDIO, space_before=6)

# ── Slide 3 — O que é ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "O que é o NXProject")
add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), AZUL_MEDIO)
add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(1.2),
            "Uma camada de planejamento sobre o Azure DevOps que transforma o backlog técnico "
            "em cronograma gerenciável — com datas, dependências, Gantt, alocação de recursos e "
            "rastreabilidade bidirecional — sem mudar o fluxo de trabalho da equipe técnica.",
            14, color=BRANCO, align=PP_ALIGN.CENTER)
pilares = [
    ("🗓️", "Cronograma\nautomático",
     "Datas calculadas a partir de duração, calendário de trabalho, feriados e predecessoras."),
    ("🔗", "Integração\nbidirecional",
     "Importa e sincroniza com Azure DevOps: work items, sprints, responsáveis, estimativas."),
    ("📊", "Gantt\ninterativo",
     "Visualização de barras, dependências, marcos e linha do tempo com zoom ajustável."),
    ("👥", "Gestão de\nrecursos",
     "Visão de alocação por pessoa, identificação de sobrecargas e conflitos de agenda."),
    ("🤖", "Assistente\nIA",
     "Sugere estrutura de tarefas e decomposição de histórias a partir de texto livre."),
]
for i, (emoji, titulo, corpo) in enumerate(pilares):
    l = Inches(0.3 + i * 2.56); t = Inches(3.25); w = Inches(2.4); h = Inches(3.8)
    add_rect(slide, l, t, w, h, BRANCO)
    add_textbox(slide, l, t+Inches(0.2), w, Inches(0.5), emoji, 26, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t+Inches(0.75), w-Inches(0.2), Inches(0.75),
                titulo, 13, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.15), t+Inches(1.6), w-Inches(0.3), Inches(1.9),
                corpo, 10.5, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# ── Slide 4 — Funcionalidades ─────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Principais funcionalidades")
funcs = [
    ("Importação Azure DevOps",
     "• Hierarquia Project → Epic → Feature → Story\n• Estimativas, datas, sprints e responsáveis\n• Predecessoras e bloqueios (tag Block)"),
    ("Cronograma inteligente",
     "• Cascata automática por predecessoras\n• Predecessora virtual por recurso\n• Calendário com feriados e horas úteis"),
    ("Gráfico Gantt",
     "• Barras, marcos (milestones), setas de dependência\n• Zoom: Dia, Semana, Sprint, Mês, Trimestre\n• Arrastar para replanejar visualmente"),
    ("Gestão de recursos",
     "• Alocação por pessoa e sprint\n• Alerta de sobrecarga (>100%)\n• Filtro de cronograma por pessoa"),
    ("Sincronização bidirecional",
     "• Datas, horas, estado, sprint, tags e predecessoras\n• Cria work items novos no DevOps\n• Relatório detalhado de sincronização"),
    ("Health Check",
     "• Atividades em atraso (Fim < hoje e % < 100)\n• Itens sem responsável\n• Dependências circulares e bloqueios"),
]
for i, (titulo, corpo) in enumerate(funcs):
    col = i % 3; row = i // 3
    l = Inches(0.4 + col * 4.3); t = Inches(1.55 + row * 2.85); w = Inches(4.0); h = Inches(2.6)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, w, Inches(0.45), AZUL_MEDIO)
    add_textbox(slide, l+Inches(0.15), t+Inches(0.05), w-Inches(0.2), Inches(0.38),
                titulo, 12, bold=True, color=BRANCO)
    add_textbox(slide, l+Inches(0.15), t+Inches(0.52), w-Inches(0.25), h-Inches(0.6),
                corpo, 10.5, color=CINZA_MEDIO)

# ── Slide 5 — Vantagens para a Gestão ────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Vantagens estratégicas para a Gestão")
vantagens = [
    ("🎯  Visibilidade real sobre os projetos de TI",
     "Cronograma consolidado com datas, dependências e percentual de conclusão — sem depender de relatórios manuais da equipe."),
    ("⚡  Zero impacto no fluxo técnico",
     "A equipe continua no Azure DevOps como antes. O NXProject é uma camada de leitura e planejamento — nenhuma mudança de processo necessária."),
    ("💰  Redução de retrabalho e custos ocultos",
     "Elimina planilhas paralelas, reuniões de status demoradas e o custo de atrasos não detectados a tempo."),
    ("📈  Decisão baseada em dados",
     "Health Check automático, alocação de recursos e alertas proativos de atraso dão à gestão informação precisa para tomar decisões rápidas."),
    ("🔒  Segurança e controle",
     "Credenciais Azure DevOps protegidas por DPAPI (criptografia Windows por usuário). Dados do projeto em arquivo local — sem cloud obrigatório."),
    ("📦  Implantação simples",
     "Instalação standalone (zero dependências externas), pronto para uso em minutos. Sem infraestrutura adicional de servidor."),
]
for i, (titulo, corpo) in enumerate(vantagens):
    col = i % 2; row = i // 2
    l = Inches(0.4 + col * 6.4); t = Inches(1.55 + row * 1.9); w = Inches(6.1); h = Inches(1.75)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, Inches(0.07), h, LARANJA)
    tb = slide.shapes.add_textbox(l+Inches(0.22), t+Inches(0.12), w-Inches(0.35), h-Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run(); run.text = titulo
    run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = AZUL_ESCURO
    add_para(tf, corpo, 10.5, color=CINZA_MEDIO, space_before=5)

# ── Slide 6 — Fluxo ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Como funciona — fluxo de trabalho")
etapas = [
    ("1", "Importar",    "Conecte ao Azure DevOps\ne importe a hierarquia\nde work items"),
    ("2", "Planejar",    "Ajuste datas, duração,\nrecursos e dependências\nno cronograma"),
    ("3", "Visualizar",  "Acompanhe no Gantt\ndatas, marcos e\nalocação de recursos"),
    ("4", "Sincronizar", "Envie de volta ao DevOps\ndatas, horas, estados\ne novos work items"),
    ("5", "Monitorar",   "Health Check automático\nalerta atrasos,\nbloqueios e sobrecarga"),
]
step_w = Inches(2.2); arrow_w = Inches(0.4)
total_w = len(etapas) * step_w + (len(etapas) - 1) * arrow_w
start_l = (W - total_w) / 2
t_box = Inches(2.2); h_box = Inches(3.8)
for i, (num, titulo, corpo) in enumerate(etapas):
    l = start_l + i * (step_w + arrow_w)
    add_rect(slide, l, t_box, step_w, h_box, AZUL_MEDIO)
    add_rect(slide, l+Inches(0.75), t_box+Inches(0.18), Inches(0.7), Inches(0.7), LARANJA)
    add_textbox(slide, l+Inches(0.75), t_box+Inches(0.18), Inches(0.7), Inches(0.7),
                num, 20, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t_box+Inches(1.1), step_w-Inches(0.2), Inches(0.6),
                titulo, 14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t_box+Inches(1.85), step_w-Inches(0.2), Inches(1.8),
                corpo, 11, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
    if i < len(etapas) - 1:
        ax = l + step_w + Inches(0.1); ay = t_box + h_box/2 - Inches(0.15)
        add_textbox(slide, ax, ay, arrow_w-Inches(0.1), Inches(0.3),
                    "→", 22, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)

# ── Slide 7 — Para quem é ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Para quem é o NXProject")
perfis = [
    ("🧑‍💼", "Gerente\nde Projeto",
     "Cronograma integrado ao backlog, alertas de atraso, visão de dependências e datas negociadas."),
    ("🔄", "Scrum Master\n/ RTE",
     "Capacidade por sprint, conflito de alocação, impacto de mudanças e cascata automática."),
    ("💻", "Tech Lead",
     "Visão de Features e Stories com predecessoras, estimativas em horas e rastreabilidade."),
    ("📊", "Gestão\n/ PMO",
     "Visão consolidada do portfólio, exportação para MS Project/Excel e Health Check executivo."),
]
for i, (emoji, titulo, corpo) in enumerate(perfis):
    l = Inches(0.5 + i * 3.2); t = Inches(1.7); w = Inches(2.9); h = Inches(5.1)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, w, Inches(0.06), LARANJA)
    add_textbox(slide, l, t+Inches(0.2), w, Inches(0.7), emoji, 34, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t+Inches(1.0), w-Inches(0.2), Inches(0.8),
                titulo, 14, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.15), t+Inches(1.9), w-Inches(0.3), Inches(3.0),
                corpo, 11, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# ── Slide 8 — Encerramento ────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, AZUL_ESCURO)
add_rect(slide, 0, Inches(3.5), W, Inches(0.07), LARANJA)
add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
            "Transforme seu backlog DevOps em um cronograma gerenciável.",
            26, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.1),
            "O NXProject conecta a gestão ao desenvolvimento — com transparência, "
            "automação e rastreabilidade — sem mudar o processo técnico.",
            16, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.7),
            "Solicite uma demonstração", 22, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.75), Inches(11.3), Inches(0.5),
            "comercial.nexus.xdata@gmail.com  •  nexusxdata.com.br",
            14, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
            "Nexus XData Tecnologia Ltda",
            11, color=RGBColor(0x7A, 0x9A, 0xC8), align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "NXProject_Gestao_Inteligente_DevOps.pptx")
prs.save(out)
print(f"PPT salvo em: {out}")
